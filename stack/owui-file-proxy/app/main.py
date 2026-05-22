import os
import glob
import time
import uuid
import re
import hmac
import hashlib
import base64
import binascii
import mimetypes
import json
import sqlite3
import html
import zipfile
import io
import copy
import requests
from contextvars import ContextVar
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional, Any

from fastapi import FastAPI, HTTPException, Header, Query, Request
from fastapi.responses import Response, JSONResponse, FileResponse, RedirectResponse
from pydantic import BaseModel, Field, constr

try:
    from docx import Document  # type: ignore
    from docx.shared import RGBColor  # type: ignore
except Exception:  # pragma: no cover - optional runtime dependency
    Document = None  # type: ignore
    RGBColor = None  # type: ignore

try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore
except Exception:  # pragma: no cover - optional runtime dependency
    Presentation = None  # type: ignore
    Inches = None  # type: ignore
    Pt = None  # type: ignore

try:
    from reportlab.lib import colors  # type: ignore
    from reportlab.lib.pagesizes import A4  # type: ignore
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet  # type: ignore
    from reportlab.lib.units import mm  # type: ignore
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle  # type: ignore
except Exception:  # pragma: no cover - optional runtime dependency
    colors = None  # type: ignore
    A4 = None  # type: ignore
    ParagraphStyle = None  # type: ignore
    getSampleStyleSheet = None  # type: ignore
    mm = None  # type: ignore
    Paragraph = None  # type: ignore
    SimpleDocTemplate = None  # type: ignore
    Spacer = None  # type: ignore
    Table = None  # type: ignore
    TableStyle = None  # type: ignore

try:
    from pypdf import PdfReader, PdfWriter  # type: ignore
except Exception:  # pragma: no cover - optional runtime dependency
    PdfReader = None  # type: ignore
    PdfWriter = None  # type: ignore

OWUI_UPLOAD_DIR = os.getenv("OWUI_UPLOAD_DIR", "/app/backend/data/uploads")
DOC_WORKER_URL = os.getenv("DOC_WORKER_URL", "http://document-worker:8090")
KAHLE_ASSETS_ROOT = Path(os.getenv("KAHLE_ASSETS_ROOT", "/assets")).resolve()
KAHLE_DOCX_TEMPLATE = Path(os.getenv("KAHLE_DOCX_TEMPLATE", str(KAHLE_ASSETS_ROOT / "templates/docx/KAHLE-DOCX-VORLAGE.docx")))
KAHLE_PDF_TEMPLATE = Path(os.getenv("KAHLE_PDF_TEMPLATE", str(KAHLE_ASSETS_ROOT / "templates/pdf/KAHLE-PDF-VORLAGE.pdf")))
KAHLE_PPTX_TEMPLATE = Path(os.getenv("KAHLE_PPTX_TEMPLATE", str(KAHLE_ASSETS_ROOT / "templates/pptx/KAHLE-PPTX-Vorlage.pptx")))
KAHLE_BRAND_CONFIG = Path(os.getenv("KAHLE_BRAND_CONFIG", str(KAHLE_ASSETS_ROOT / "brand/colors/kahle-brand.json")))
KAHLE_LOGO_PRIMARY = Path(
    os.getenv("KAHLE_LOGO_PRIMARY", str(KAHLE_ASSETS_ROOT / "brand/logos/Logo_Kahle_Gruppe_positiv.png"))
)

# Backwards-compatible: TOOL_API_KEY protects proxy endpoints (optional),
# and is also forwarded to worker as x-api-key (optional).
TOOL_API_KEY = os.getenv("TOOL_API_KEY", "")
WORKER_API_KEY = os.getenv("WORKER_API_KEY", "") or TOOL_API_KEY
REQUIRE_TOOL_API_KEY = os.getenv("REQUIRE_TOOL_API_KEY", "true").lower() == "true"

UPLOAD_ROOT = Path(OWUI_UPLOAD_DIR).resolve()

# Output subfolder for saved files
OUTPUT_SUBDIR = os.getenv("OWUI_OUTPUT_SUBDIR", "edited").strip("/")
OWUI_DB_PATH = Path(os.getenv("OWUI_DB_PATH", "/app/backend/data/webui.db"))
OWUI_LOG_CLEANUP_REPORT_PATH = Path(
    os.getenv("OWUI_LOG_CLEANUP_REPORT_PATH", "/retention-reports/openwebui_log_cleanup_report.json")
)
OWUI_LOG_CLEANUP_STALE_HOURS = int(os.getenv("OWUI_LOG_CLEANUP_STALE_HOURS", "36"))

# Size limit
MAX_FILE_BYTES = int(os.getenv("MAX_FILE_BYTES", "20000000"))  # 20 MB

# Signed download links
FILE_LINK_SECRET = os.getenv("FILE_LINK_SECRET", "")
FILE_LINK_TTL_SECONDS = int(os.getenv("FILE_LINK_TTL_SECONDS", "3600"))
ALLOW_UNSIGNED_DOWNLOADS = os.getenv("ALLOW_UNSIGNED_DOWNLOADS", "false").lower() == "true"

# Optional: absolute links in chat
PUBLIC_BASE_URL = os.getenv("PUBLIC_BASE_URL", "").rstrip("/")

# Multi-user safety toggles
REQUIRE_EXACT_FILE_PATH = os.getenv("REQUIRE_EXACT_FILE_PATH", "false").lower() == "true"
DISALLOW_WILDCARDS = os.getenv("DISALLOW_WILDCARDS", "true").lower() == "true"
RECENT_UPLOAD_DISAMBIGUATION_SECONDS = int(os.getenv("RECENT_UPLOAD_DISAMBIGUATION_SECONDS", "3600"))

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PDF_MIME = "application/pdf"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

# Allow saving these file types from worker results
ALLOWED_SAVE_EXT = {".docx", ".md", ".txt", ".csv", ".pdf", ".xlsx", ".pptx"}

WILDCARD_RE = re.compile(r"[*?\[\]{}]")
OWUI_UUID_PREFIX_RE = re.compile(r"^[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}_.+", re.I)
NonEmptyStr = constr(strip_whitespace=True, min_length=1)


app = FastAPI(title="OWUI File Proxy", version="1.5.0")

_REQUEST_API_KEY: ContextVar[str] = ContextVar("request_api_key", default="")
_DOCX_TEMPLATE_USED: ContextVar[bool] = ContextVar("docx_template_used", default=False)
_PDF_TEMPLATE_USED: ContextVar[bool] = ContextVar("pdf_template_used", default=False)
_PPTX_TEMPLATE_USED: ContextVar[bool] = ContextVar("pptx_template_used", default=False)
AUTH_EXEMPT_PATHS = {"/health", "/files/download"}


# -----------------------------
# Security helpers
# -----------------------------
def _extract_bearer_token(authorization: Optional[str]) -> str:
    if not authorization:
        return ""
    scheme, _, token = authorization.strip().partition(" ")
    if scheme.lower() != "bearer" or not token.strip():
        return ""
    return token.strip()


def _provided_api_key(x_api_key: Optional[str], authorization: Optional[str]) -> str:
    return (x_api_key or "").strip() or _extract_bearer_token(authorization)


def _validate_api_key(provided: str) -> bool:
    return bool(TOOL_API_KEY) and bool(provided) and hmac.compare_digest(provided, TOOL_API_KEY)


def _require_api_key(x_api_key: Optional[str] = None, authorization: Optional[str] = None) -> None:
    if not TOOL_API_KEY:
        if REQUIRE_TOOL_API_KEY:
            raise HTTPException(status_code=500, detail="TOOL_API_KEY not configured")
        return

    provided = _provided_api_key(x_api_key, authorization) or _REQUEST_API_KEY.get()
    if not _validate_api_key(provided):
        raise HTTPException(status_code=401, detail="unauthorized")


@app.middleware("http")
async def enforce_api_key(request: Request, call_next):
    path = request.url.path
    if path not in AUTH_EXEMPT_PATHS:
        provided = _provided_api_key(
            request.headers.get("x-api-key"),
            request.headers.get("authorization"),
        )
        if TOOL_API_KEY:
            if not _validate_api_key(provided):
                return JSONResponse(status_code=401, content={"detail": "unauthorized"})
        elif REQUIRE_TOOL_API_KEY:
            return JSONResponse(status_code=500, content={"detail": "TOOL_API_KEY not configured"})

        token = _REQUEST_API_KEY.set(provided)
        try:
            return await call_next(request)
        finally:
            _REQUEST_API_KEY.reset(token)

    return await call_next(request)


def _ensure_within_upload_root(resolved: Path) -> None:
    try:
        resolved.relative_to(UPLOAD_ROOT)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path: outside uploads directory.")


def _safe_relpath(rel: str) -> str:
    rel = (rel or "").replace("\\", "/").lstrip("/")
    parts = [p for p in rel.split("/") if p]
    if any(p == ".." for p in parts):
        raise HTTPException(status_code=400, detail="invalid_path")
    return "/".join(parts)


def _decode_literal_unicode_escapes(value: str) -> str:
    def replace(match: re.Match[str]) -> str:
        try:
            return chr(int(match.group(1), 16))
        except Exception:
            return match.group(0)

    return re.sub(r"(?:\\+u|_u)([0-9a-fA-F]{4})", replace, str(value or ""))


def _ascii_filename_text(value: str) -> str:
    text = _decode_literal_unicode_escapes(value)
    replacements = {
        "ä": "ae",
        "ö": "oe",
        "ü": "ue",
        "ß": "ss",
        "Ä": "Ae",
        "Ö": "Oe",
        "Ü": "Ue",
    }
    for source, target in replacements.items():
        text = text.replace(source, target)
    return text


def _sanitize_filename(name: str) -> str:
    name = _decode_literal_unicode_escapes(name)
    name = re.split(r"[\\/]", _ascii_filename_text(name))[-1]
    name = re.sub(r"[^A-Za-z0-9._-]+", "_", name).strip("._")
    if not name:
        name = f"file_{uuid.uuid4().hex}"
    return name


def _reject_wildcards(raw: str) -> None:
    if not raw:
        return
    if DISALLOW_WILDCARDS and WILDCARD_RE.search(raw):
        raise HTTPException(status_code=400, detail="wildcards_not_allowed_use_exact_filename")


# -----------------------------
# File selection
# -----------------------------
def _list_files(pattern: str) -> list[Path]:
    files = [Path(p) for p in glob.glob(str(UPLOAD_ROOT / pattern), recursive=True)]
    return [p for p in files if p.is_file()]


def _is_openwebui_upload_for_name(path: Path, visible_name: str) -> bool:
    return path.name.endswith(f"_{visible_name}") and bool(OWUI_UUID_PREFIX_RE.match(path.name))


def _pick_recent_upload_duplicate(hits: list[Path], visible_name: str) -> Path | None:
    if not hits or RECENT_UPLOAD_DISAMBIGUATION_SECONDS <= 0:
        return None

    now = time.time()
    upload_hits = [h for h in hits if _is_openwebui_upload_for_name(h, visible_name)]
    if not upload_hits:
        return None

    upload_hits.sort(key=lambda p: p.stat().st_mtime, reverse=True)
    newest = upload_hits[0]
    if now - newest.stat().st_mtime > RECENT_UPLOAD_DISAMBIGUATION_SECONDS:
        return None

    return newest


def _pick_latest_file_by_exts(allowed_exts: tuple[str, ...]) -> Path:
    # NOTE: In strict mode we generally don't want "latest". This exists for backwards compatibility.
    candidates: list[Path] = []
    for ext in allowed_exts:
        ext = ext.lstrip(".")
        candidates += _list_files(f"**/*.{ext}")

    if not candidates:
        raise HTTPException(
            status_code=400,
            detail=f"No matching files found in OWUI_UPLOAD_DIR={UPLOAD_ROOT} for extensions={allowed_exts}",
        )

    candidates.sort(key=lambda p: p.stat().st_mtime, reverse=True)
    return candidates[0]


def _resolve_path(
    file_path: str | None,
    allowed_exts: tuple[str, ...],
    require_exact: bool = False,
    preferred_exact_paths: Optional[list[str]] = None,
) -> Path:
    """
    Resolve file_path safely inside UPLOAD_ROOT.

    Accepts:
      - filename only (e.g. 'test.docx') -> search recursively in uploads
      - relative path under uploads -> resolve
      - absolute path, but must still be inside uploads

    Strict mode / exact mode:
      - file_path MUST be provided and non-null.
      - ambiguous bare filenames are rejected; user must provide the exact relative path.
    """
    strict = REQUIRE_EXACT_FILE_PATH or require_exact
    if strict and (file_path is None or str(file_path).strip() in ("", "null", "None")):
        raise HTTPException(status_code=400, detail="file_path_required_use_exact_uploaded_filename")

    if not file_path:
        return _pick_latest_file_by_exts(allowed_exts)

    raw = file_path.strip()
    if raw in ("null", "None", ""):
        return _pick_latest_file_by_exts(allowed_exts)

    # Normalize common client-side path hints.
    raw = raw.replace("\\", "/").lstrip("/")
    if raw.startswith("./"):
        raw = raw[2:]
    if raw.startswith("uploads/"):
        raw = raw[len("uploads/") :]

    # Reject obvious filename placeholders generated by model guesses.
    lowered = raw.lower()
    if any(tok in lowered for tok in ("dateiname", "datenname", "your_file", "your_", "anhang", "<", ">")):
        raise HTTPException(status_code=400, detail="placeholder_filename_not_allowed_use_exact_uploaded_filename")

    _reject_wildcards(raw)

    p = Path(raw)

    # filename only -> search
    if ("/" not in raw) and ("\\" not in raw):
        hits = _list_files(f"**/{raw}")
        if not hits:
            # OWUI often stores files as "<uuid>_<original_name>".
            hits = _list_files(f"**/*_{raw}")
        if not hits:
            raise HTTPException(status_code=400, detail=f"File not found in uploads: {raw}")

        # Optional deterministic disambiguation using exact paths from the current message attachments.
        if len(hits) > 1 and preferred_exact_paths:
            preferred_set = set()
            for p in preferred_exact_paths:
                if not isinstance(p, str):
                    continue
                cand = p.strip().replace("\\", "/").lstrip("/")
                if cand.startswith("uploads/"):
                    cand = cand[len("uploads/") :]
                if cand:
                    preferred_set.add(cand)
                    preferred_set.add(cand.split("/")[-1])

            filtered = []
            for h in hits:
                try:
                    rel = str(h.resolve().relative_to(UPLOAD_ROOT)).replace("\\", "/")
                except Exception:
                    rel = h.name
                if rel in preferred_set or h.name in preferred_set:
                    filtered.append(h)

            if len(filtered) == 1:
                hits = filtered

        if strict and len(hits) > 1:
            recent_upload = _pick_recent_upload_duplicate(hits, raw)
            if recent_upload is not None:
                hits = [recent_upload]

        if strict and len(hits) > 1:
            matches = []
            for h in hits[:20]:
                try:
                    rel = str(h.resolve().relative_to(UPLOAD_ROOT)).replace("\\", "/")
                except Exception:
                    rel = h.name
                matches.append(rel)
            raise HTTPException(
                status_code=400,
                detail={
                    "error": "ambiguous_filename_use_exact_relative_path",
                    "filename": raw,
                    "matches": matches,
                },
            )
        hits.sort(key=lambda x: x.stat().st_mtime, reverse=True)
        resolved = hits[0].resolve()
    else:
        # treat as path
        if p.is_absolute():
            resolved = p.resolve()
        else:
            resolved = (UPLOAD_ROOT / p).resolve()

    _ensure_within_upload_root(resolved)

    allowed_lower = [e.lower() for e in allowed_exts]
    if resolved.suffix.lower() not in allowed_lower:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid file type: {resolved.suffix}. Allowed: {', '.join(allowed_exts)}",
        )

    if not resolved.is_file():
        raise HTTPException(status_code=400, detail=f"File not found: {resolved}")

    return resolved


def _read_file(path: Path) -> bytes:
    size = path.stat().st_size
    if size > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail=f"file_too_large (>{MAX_FILE_BYTES} bytes)")
    with open(path, "rb") as f:
        return f.read()


# -----------------------------
# Worker headers (forwarding)
# -----------------------------
def _worker_headers() -> dict:
    h = {}
    if WORKER_API_KEY:
        h["x-api-key"] = WORKER_API_KEY
    return h


def _raise_for_worker_response(r: requests.Response) -> None:
    """
    Preserve upstream worker error semantics for tool callers.
    - 4xx from worker should stay 4xx (validation/user-fixable)
    - 5xx from worker is mapped to 502
    """
    if r.status_code == 200:
        return

    detail: Any
    try:
        payload = r.json()
        detail = payload.get("detail", payload)
    except Exception:
        detail = (r.text or "").strip()[:2000] or f"worker_http_error_{r.status_code}"

    status = r.status_code if 400 <= r.status_code < 500 else 502
    raise HTTPException(status_code=status, detail=detail)


# -----------------------------
# Signing + saving
# -----------------------------
def _sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _b64decode_strict(b64: str) -> bytes:
    try:
        data = base64.b64decode(b64, validate=True)
    except binascii.Error:
        raise HTTPException(status_code=400, detail="invalid_base64")
    if len(data) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail=f"file_too_large (>{MAX_FILE_BYTES} bytes)")
    return data


def _write_atomic(path: Path, data: bytes) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp = path.with_suffix(path.suffix + ".tmp")
    with open(tmp, "wb") as f:
        f.write(data)
    os.replace(tmp, path)


def _file_exists(path: Path) -> bool:
    try:
        return path.exists() and path.is_file()
    except Exception:
        return False


def _load_brand_config() -> dict[str, Any]:
    if not _file_exists(KAHLE_BRAND_CONFIG):
        return {}
    try:
        with open(KAHLE_BRAND_CONFIG, "r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, dict) else {}
    except Exception:
        return {}


def _brand_value(section: str, key: str, default: str) -> str:
    data = _load_brand_config()
    section_data = data.get(section, {})
    if isinstance(section_data, dict):
        value = section_data.get(key)
        if isinstance(value, str) and value.strip():
            return value.strip().lstrip("#")
    return default


def _safe_style(document: Any, preferred: str, fallback: str = "Normal") -> str:
    try:
        _ = document.styles[preferred]
        return preferred
    except Exception:
        return fallback


def _strip_single_markdown_markers(text: str) -> str:
    cleaned = re.sub(r"(?<!\*)\*([^*\n]+)\*(?!\*)", r"\1", text or "")
    cleaned = re.sub(r"(?<!\w)_([^_\n]+)_(?!\w)", r"\1", cleaned)
    return cleaned.replace("**", "").replace("__", "")


def _markdown_inline_segments(text: str, bold: bool = False) -> list[tuple[str, bool]]:
    segments: list[tuple[str, bool]] = []
    pos = 0
    for match in re.finditer(r"(\*\*|__)(.+?)\1", text or ""):
        if match.start() > pos:
            plain = _strip_single_markdown_markers((text or "")[pos : match.start()])
            if plain:
                segments.append((plain, bool(bold)))
        strong = _strip_single_markdown_markers(match.group(2))
        if strong:
            segments.append((strong, True))
        pos = match.end()

    tail = _strip_single_markdown_markers((text or "")[pos:])
    if tail or not segments:
        segments.append((tail, bool(bold)))
    return segments


def _add_docx_paragraph(document: Any, text: str, style: str = "Normal", bold: bool = False, color_hex: str = "") -> None:
    paragraph = document.add_paragraph(style=_safe_style(document, style))
    for segment_text, segment_bold in _markdown_inline_segments(text or "", bold):
        run = paragraph.add_run(segment_text)
        run.bold = bool(segment_bold)
        if color_hex and RGBColor is not None:
            value = color_hex.strip().lstrip("#")
            try:
                run.font.color.rgb = RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
            except Exception:
                pass


def _markdown_lines(content: str) -> list[str]:
    lines: list[str] = []
    in_code = False
    for raw_line in (content or "").splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            lines.append(line)
        else:
            lines.append(stripped)
    return lines


def _docx_run_xml(text: str, bold: bool = False) -> str:
    escaped = html.escape(text or "")
    run_props = "<w:rPr><w:b/></w:rPr>" if bold else ""
    return f'<w:r>{run_props}<w:t xml:space="preserve">{escaped}</w:t></w:r>'


def _docx_paragraph_xml(text: str, style: str = "", bold: bool = False) -> str:
    style_xml = f'<w:pPr><w:pStyle w:val="{style}"/></w:pPr>' if style else ""
    runs = "".join(_docx_run_xml(segment_text, segment_bold) for segment_text, segment_bold in _markdown_inline_segments(text or "", bold))
    return f"<w:p>{style_xml}{runs}</w:p>"


def _generated_at_label() -> str:
    return time.strftime("%Y-%m-%d %H:%M Europe/Berlin")


def _markdown_to_template_docx_bytes(content: str, title: str = "Dokument") -> Optional[bytes]:
    if Document is None or not _file_exists(KAHLE_DOCX_TEMPLATE):
        return None

    try:
        document = Document(str(KAHLE_DOCX_TEMPLATE))

        # Use the supplied file as style/theme template, but replace body content.
        body = document._element.body  # noqa: SLF001 - python-docx has no public clear-body API.
        section_props = None
        for child in list(body):
            if child.tag.endswith("sectPr"):
                section_props = child
                continue
            body.remove(child)
        if section_props is not None and section_props.getparent() is None:
            body.append(section_props)

        clean_title = (title or "Dokument").strip() or "Dokument"
        brand_blue = _brand_value("colors", "blue", "0069B3")
        brand_ink = _brand_value("colors", "ink", "0F2430")
        _add_docx_paragraph(document, clean_title, "Title", bold=True, color_hex=brand_blue)
        _add_docx_paragraph(document, f"Erstellt mit KAHLE-Vinci | Stand: {_generated_at_label()}", "Subtitle")
        document.add_paragraph()

        first_content_heading_seen = False
        for stripped in _markdown_lines(content):
            if not stripped:
                document.add_paragraph()
                continue
            if stripped.startswith("#"):
                heading = stripped.lstrip("#").strip()
                if not first_content_heading_seen and heading.lower() == clean_title.lower():
                    first_content_heading_seen = True
                    continue
                first_content_heading_seen = True
                level = min(max(len(stripped) - len(stripped.lstrip("#")), 1), 3)
                _add_docx_paragraph(
                    document,
                    heading,
                    f"Heading {level}",
                    bold=True,
                    color_hex=brand_blue if level == 1 else brand_ink,
                )
                continue
            if stripped.startswith(("- ", "* ")):
                _add_docx_paragraph(document, stripped[2:].strip(), "List Bullet")
                continue
            numbered = re.match(r"^\d+[.)]\s+(.+)$", stripped)
            if numbered:
                _add_docx_paragraph(document, stripped, "List Number")
                continue
            _add_docx_paragraph(document, stripped)

        out = io.BytesIO()
        document.save(out)
        _DOCX_TEMPLATE_USED.set(True)
        return out.getvalue()
    except Exception:
        return None


def _markdown_to_docx_bytes(content: str, title: str = "Dokument") -> bytes:
    """
    Render plain Markdown-ish text into a simple valid DOCX using only stdlib.
    This intentionally favors dependable downloadable files over advanced layout.
    """
    templated = _markdown_to_template_docx_bytes(content, title)
    if templated:
        return templated

    paragraphs: list[str] = []
    title = (title or "Dokument").strip() or "Dokument"
    paragraphs.append(_docx_paragraph_xml(title, "Title", bold=True))
    paragraphs.append(_docx_paragraph_xml(f"Erstellt mit KAHLE-Vinci | Stand: {_generated_at_label()}", "Subtitle"))
    paragraphs.append(_docx_paragraph_xml(""))

    in_code = False
    first_content_heading_seen = False
    for raw_line in (content or "").splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if not stripped:
            paragraphs.append(_docx_paragraph_xml(""))
            continue
        if in_code:
            paragraphs.append(_docx_paragraph_xml(line, "Code"))
            continue
        if stripped.startswith("#"):
            heading = stripped.lstrip("#").strip()
            if not first_content_heading_seen and heading.lower() == title.lower():
                first_content_heading_seen = True
                continue
            first_content_heading_seen = True
            level = min(max(len(stripped) - len(stripped.lstrip("#")), 1), 3)
            paragraphs.append(_docx_paragraph_xml(heading, f"Heading{level}", bold=True))
            continue
        if stripped.startswith(("- ", "* ")):
            paragraphs.append(_docx_paragraph_xml(f"- {stripped[2:].strip()}"))
            continue
        numbered = re.match(r"^\d+[.)]\s+(.+)$", stripped)
        if numbered:
            paragraphs.append(_docx_paragraph_xml(stripped))
            continue
        paragraphs.append(_docx_paragraph_xml(stripped))

    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body>"
        f"{''.join(paragraphs)}"
        '<w:sectPr><w:pgSz w:w="11906" w:h="16838"/><w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"/></w:sectPr>'
        "</w:body></w:document>"
    )
    styles_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        '<w:docDefaults><w:rPrDefault><w:rPr><w:rFonts w:ascii="Arial" w:hAnsi="Arial"/><w:sz w:val="22"/></w:rPr></w:rPrDefault></w:docDefaults>'
        '<w:style w:type="paragraph" w:default="1" w:styleId="Normal"><w:name w:val="Normal"/><w:pPr><w:spacing w:after="120" w:line="276" w:lineRule="auto"/></w:pPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Title"><w:name w:val="Title"/><w:pPr><w:spacing w:after="80"/></w:pPr><w:rPr><w:b/><w:sz w:val="40"/><w:color w:val="0069B3"/></w:rPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Subtitle"><w:name w:val="Subtitle"/><w:pPr><w:spacing w:after="240"/></w:pPr><w:rPr><w:sz w:val="18"/><w:color w:val="666666"/></w:rPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Heading1"><w:name w:val="heading 1"/><w:pPr><w:spacing w:before="260" w:after="120"/></w:pPr><w:rPr><w:b/><w:sz w:val="30"/><w:color w:val="0069B3"/></w:rPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Heading2"><w:name w:val="heading 2"/><w:pPr><w:spacing w:before="200" w:after="80"/></w:pPr><w:rPr><w:b/><w:sz w:val="26"/><w:color w:val="0F2430"/></w:rPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Heading3"><w:name w:val="heading 3"/><w:pPr><w:spacing w:before="160" w:after="60"/></w:pPr><w:rPr><w:b/><w:sz w:val="23"/><w:color w:val="0F2430"/></w:rPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Code"><w:name w:val="Code"/><w:rPr><w:rFonts w:ascii="Consolas" w:hAnsi="Consolas"/><w:sz w:val="19"/></w:rPr></w:style>'
        "</w:styles>"
    )
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        '<Override PartName="/word/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
        "</Relationships>"
    )
    doc_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>'
        "</Relationships>"
    )

    import io

    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("word/document.xml", document_xml)
        zf.writestr("word/styles.xml", styles_xml)
        zf.writestr("word/_rels/document.xml.rels", doc_rels)
    return out.getvalue()


def _markdown_to_plain_lines(content: str, title: str = "Dokument") -> list[str]:
    clean_title = (title or "Dokument").strip() or "Dokument"
    lines = [clean_title, f"Erstellt mit KAHLE-Vinci | Stand: {_generated_at_label()}", ""]
    in_code = False
    first_content_heading_seen = False
    for raw_line in (content or "").splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if not stripped:
            lines.append("")
            continue
        if in_code:
            lines.append(line)
            continue
        if stripped.startswith("#"):
            heading = stripped.lstrip("#").strip()
            if not first_content_heading_seen and heading.lower() == clean_title.lower():
                first_content_heading_seen = True
                continue
            first_content_heading_seen = True
            lines.append(heading)
            continue
        if stripped.startswith(("- ", "* ")):
            lines.append(f"- {stripped[2:].strip()}")
            continue
        lines.append(stripped)
    return lines


def _wrap_pdf_lines(lines: list[str], width: int = 92) -> list[str]:
    wrapped: list[str] = []
    for line in lines:
        if not line:
            wrapped.append("")
            continue
        words = line.split()
        current = ""
        for word in words:
            candidate = f"{current} {word}".strip()
            if len(candidate) <= width:
                current = candidate
                continue
            if current:
                wrapped.append(current)
            current = word
        if current:
            wrapped.append(current)
    return wrapped


def _pdf_escape_text(text: str) -> str:
    encoded = (text or "").encode("cp1252", errors="replace").decode("cp1252")
    return encoded.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def _markdown_inline_to_reportlab(text: str) -> str:
    safe = html.escape(text or "")
    safe = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", safe)
    safe = re.sub(r"\*(.+?)\*", r"<i>\1</i>", safe)
    return safe


def _text_to_reportlab_pdf_bytes(content: str, title: str = "Dokument") -> Optional[bytes]:
    if not all([colors, A4, ParagraphStyle, getSampleStyleSheet, mm, Paragraph, SimpleDocTemplate, Spacer]):
        return None

    try:
        clean_title = (title or "Dokument").strip() or "Dokument"
        brand_ink = _brand_value("colors", "ink", "0F2430")
        brand_blue = _brand_value("colors", "blue", "005A8F")
        brand_muted = _brand_value("colors", "muted", "6B7280")
        font_body = _brand_value("fonts", "body", "Helvetica")
        font_heading = _brand_value("fonts", "headings", font_body)

        styles = getSampleStyleSheet()
        styles.add(
            ParagraphStyle(
                name="KahleTitle",
                parent=styles["Title"],
                fontName="Helvetica-Bold",
                fontSize=22,
                leading=26,
                textColor=colors.HexColor(f"#{brand_blue}"),
                spaceAfter=8,
            )
        )
        styles.add(
            ParagraphStyle(
                name="KahleSubtitle",
                parent=styles["Normal"],
                fontName="Helvetica",
                fontSize=8.5,
                leading=11,
                textColor=colors.HexColor(f"#{brand_muted}"),
                spaceAfter=16,
            )
        )
        styles.add(
            ParagraphStyle(
                name="KahleHeading1",
                parent=styles["Heading1"],
                fontName="Helvetica-Bold",
                fontSize=15,
                leading=19,
                textColor=colors.HexColor(f"#{brand_blue}"),
                spaceBefore=12,
                spaceAfter=6,
            )
        )
        styles.add(
            ParagraphStyle(
                name="KahleHeading2",
                parent=styles["Heading2"],
                fontName="Helvetica-Bold",
                fontSize=12.5,
                leading=16,
                textColor=colors.HexColor(f"#{brand_ink}"),
                spaceBefore=10,
                spaceAfter=4,
            )
        )
        styles.add(
            ParagraphStyle(
                name="KahleBody",
                parent=styles["BodyText"],
                fontName="Helvetica",
                fontSize=9.7,
                leading=13.5,
                textColor=colors.HexColor(f"#{brand_ink}"),
                spaceAfter=5,
            )
        )
        styles.add(
            ParagraphStyle(
                name="KahleBullet",
                parent=styles["KahleBody"],
                leftIndent=12,
                firstLineIndent=-8,
            )
        )

        out = io.BytesIO()
        doc = SimpleDocTemplate(
            out,
            pagesize=A4,
            rightMargin=18 * mm,
            leftMargin=18 * mm,
            topMargin=24 * mm,
            bottomMargin=18 * mm,
            title=clean_title,
            author="KAHLE-Vinci",
        )

        story: list[Any] = [
            Paragraph(_markdown_inline_to_reportlab(clean_title), styles["KahleTitle"]),
            Paragraph(f"Erstellt mit KAHLE-Vinci | Stand: {_generated_at_label()}", styles["KahleSubtitle"]),
        ]

        first_content_heading_seen = False
        for stripped in _markdown_lines(content):
            if not stripped:
                story.append(Spacer(1, 4))
                continue
            if stripped.startswith("#"):
                heading = stripped.lstrip("#").strip()
                if not first_content_heading_seen and heading.lower() == clean_title.lower():
                    first_content_heading_seen = True
                    continue
                first_content_heading_seen = True
                level = min(max(len(stripped) - len(stripped.lstrip("#")), 1), 2)
                story.append(Paragraph(_markdown_inline_to_reportlab(heading), styles[f"KahleHeading{level}"]))
                continue
            if stripped.startswith(("- ", "* ")):
                story.append(Paragraph(f"&bull; {_markdown_inline_to_reportlab(stripped[2:].strip())}", styles["KahleBullet"]))
                continue
            story.append(Paragraph(_markdown_inline_to_reportlab(stripped), styles["KahleBody"]))

        def draw_page(canvas: Any, _doc: Any) -> None:
            canvas.saveState()
            width, height = A4
            canvas.setStrokeColor(colors.HexColor(f"#{brand_blue}"))
            canvas.setLineWidth(0.6)
            canvas.line(18 * mm, height - 17 * mm, width - 18 * mm, height - 17 * mm)
            if _file_exists(KAHLE_LOGO_PRIMARY):
                try:
                    canvas.drawImage(
                        str(KAHLE_LOGO_PRIMARY),
                        width - 50 * mm,
                        height - 16 * mm,
                        width=32 * mm,
                        height=11 * mm,
                        preserveAspectRatio=True,
                        mask="auto",
                    )
                except Exception:
                    pass
            canvas.setFillColor(colors.HexColor(f"#{brand_muted}"))
            canvas.setFont("Helvetica", 7)
            canvas.drawString(18 * mm, 10 * mm, "KAHLE-Vinci")
            canvas.drawRightString(width - 18 * mm, 10 * mm, f"Seite {canvas.getPageNumber()}")
            canvas.restoreState()

        doc.build(story, onFirstPage=draw_page, onLaterPages=draw_page)
        rendered = out.getvalue()
        return rendered
    except Exception:
        return None


def _apply_pdf_template(rendered_pdf: bytes) -> Optional[bytes]:
    """
    Use the configured PDF as a page background and place generated content on top.
    This preserves the deposited KAHLE PDF layout while keeping text generation deterministic.
    """
    if PdfReader is None or PdfWriter is None:
        return None
    if not rendered_pdf or not _file_exists(KAHLE_PDF_TEMPLATE):
        return None

    try:
        template_reader = PdfReader(str(KAHLE_PDF_TEMPLATE))
        content_reader = PdfReader(io.BytesIO(rendered_pdf))
        if not template_reader.pages or not content_reader.pages:
            return None

        writer = PdfWriter()
        last_template_index = len(template_reader.pages) - 1
        for index, content_page in enumerate(content_reader.pages):
            template_page = template_reader.pages[min(index, last_template_index)]
            page = copy.copy(template_page)
            page.merge_page(content_page)
            writer.add_page(page)

        out = io.BytesIO()
        writer.write(out)
        return out.getvalue()
    except Exception:
        return None


def _text_to_pdf_bytes(content: str, title: str = "Dokument") -> bytes:
    """
    Create a small text-only PDF with stdlib primitives.
    This is intentionally simple and dependable for downloadable briefs.
    """
    branded = _text_to_reportlab_pdf_bytes(content, title)
    if branded:
        return branded

    title = (title or "Dokument").strip() or "Dokument"
    lines = _wrap_pdf_lines(_markdown_to_plain_lines(content, title), width=88)
    pages: list[list[str]] = []
    page: list[str] = []
    for idx, line in enumerate(lines):
        if len(page) >= 48:
            pages.append(page)
            page = []
        page.append(line)
    pages.append(page or [""])

    objects: list[bytes] = []

    def add_obj(body: bytes) -> int:
        objects.append(body)
        return len(objects)

    catalog_id = add_obj(b"<< /Type /Catalog /Pages 2 0 R >>")
    pages_id = add_obj(b"<< /Type /Pages /Kids [] /Count 0 >>")
    font_id = add_obj(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica /Encoding /WinAnsiEncoding >>")
    bold_font_id = add_obj(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica-Bold /Encoding /WinAnsiEncoding >>")

    page_ids: list[int] = []
    content_ids: list[int] = []
    for page_number, page_lines in enumerate(pages, start=1):
        commands = ["BT", "/F2 18 Tf", "50 790 Td", "18 TL"]
        for idx, line in enumerate(page_lines):
            if idx:
                commands.append("T*")
            if page_number == 1 and idx == 0:
                commands.append("/F2 18 Tf")
            elif page_number == 1 and idx == 1:
                commands.append("/F1 8 Tf")
            else:
                commands.append("/F1 10 Tf")
            commands.append(f"({_pdf_escape_text(line)}) Tj")
        commands.append("ET")
        commands.append("BT")
        commands.append("/F1 8 Tf")
        commands.append("50 32 Td")
        commands.append(f"({_pdf_escape_text(f'KAHLE-Vinci | Seite {page_number} von {len(pages)}')}) Tj")
        commands.append("ET")
        stream = "\n".join(commands).encode("cp1252", errors="replace")
        content_id = add_obj(b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream")
        page_id = add_obj(
            (
                f"<< /Type /Page /Parent {pages_id} 0 R /MediaBox [0 0 595 842] "
                f"/Resources << /Font << /F1 {font_id} 0 R /F2 {bold_font_id} 0 R >> >> /Contents {content_id} 0 R >>"
            ).encode("ascii")
        )
        content_ids.append(content_id)
        page_ids.append(page_id)

    objects[pages_id - 1] = (
        f"<< /Type /Pages /Kids [{' '.join(f'{pid} 0 R' for pid in page_ids)}] /Count {len(page_ids)} >>"
    ).encode("ascii")

    out = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for idx, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out.extend(f"{idx} 0 obj\n".encode("ascii"))
        out.extend(body)
        out.extend(b"\nendobj\n")

    xref_offset = len(out)
    out.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    out.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        out.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    out.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root {catalog_id} 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    return bytes(out)


def _clear_pptx_slides(prs: Any) -> None:
    try:
        slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public clear API.
        for slide_id in list(slide_id_list):
            r_id = slide_id.rId
            prs.part.drop_rel(r_id)
            slide_id_list.remove(slide_id)
    except Exception:
        pass


def _remove_pptx_slides_after(prs: Any, keep_count: int) -> None:
    try:
        slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public remove API.
        for slide_id in list(slide_id_list)[max(keep_count, 0) :]:
            prs.part.drop_rel(slide_id.rId)
            slide_id_list.remove(slide_id)
    except Exception:
        pass


def _pick_pptx_layout(prs: Any, preferred_names: tuple[str, ...] = ()) -> Any:
    for preferred in preferred_names:
        for layout in prs.slide_layouts:
            if preferred.lower() in (getattr(layout, "name", "") or "").lower():
                return layout
    for layout in prs.slide_layouts:
        try:
            if len(layout.placeholders) >= 1:
                return layout
        except Exception:
            continue
    return prs.slide_layouts[0]


MAX_PPTX_SLIDES = 5
MAX_PPTX_BULLETS_PER_SLIDE = 4
MAX_PPTX_BULLET_CHARS = 115
MAX_PPTX_TITLE_CHARS = 58


def _shorten_pptx_text(text: str, max_chars: int = MAX_PPTX_BULLET_CHARS) -> str:
    compact = re.sub(r"\s+", " ", str(text or "")).strip()
    compact = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", compact)
    if len(compact) <= max_chars:
        return compact
    return compact[: max_chars - 1].rstrip(" ,.;:-") + "..."


def _fit_pptx_bullets(bullets: list[str]) -> list[str]:
    fitted = [_shorten_pptx_text(str(b), MAX_PPTX_BULLET_CHARS) for b in bullets if str(b).strip()]
    return fitted[:MAX_PPTX_BULLETS_PER_SLIDE]


def _markdown_to_slide_specs(content: str, title: str = "Praesentation") -> list[dict[str, Any]]:
    clean_title = (title or "Praesentation").strip() or "Praesentation"
    specs: list[dict[str, Any]] = [
        {
            "title": _shorten_pptx_text(clean_title, MAX_PPTX_TITLE_CHARS),
            "bullets": [f"Erstellt mit KAHLE-Vinci | Stand: {_generated_at_label()}"],
        }
    ]
    current: Optional[dict[str, Any]] = None

    for stripped in _markdown_lines(content):
        if not stripped:
            continue
        if stripped.startswith("#"):
            heading = _shorten_pptx_text(stripped.lstrip("#").strip(), MAX_PPTX_TITLE_CHARS)
            if heading.lower() == clean_title.lower():
                continue
            current = {"title": heading, "bullets": []}
            specs.append(current)
            continue
        if current is None:
            current = {"title": "Kernaussagen", "bullets": []}
            specs.append(current)
        if stripped.startswith(("- ", "* ")):
            current["bullets"].append(stripped[2:].strip())
        else:
            current["bullets"].append(stripped)

    expanded: list[dict[str, Any]] = []
    for spec in specs:
        bullets = [str(b).strip() for b in spec.get("bullets", []) if str(b).strip()]
        if len(bullets) <= MAX_PPTX_BULLETS_PER_SLIDE:
            expanded.append({"title": spec["title"], "bullets": _fit_pptx_bullets(bullets)})
            continue
        for idx in range(0, len(bullets), MAX_PPTX_BULLETS_PER_SLIDE):
            suffix = "" if idx == 0 else f" ({idx // MAX_PPTX_BULLETS_PER_SLIDE + 1})"
            expanded.append(
                {
                    "title": _shorten_pptx_text(f"{spec['title']}{suffix}", MAX_PPTX_TITLE_CHARS),
                    "bullets": _fit_pptx_bullets(bullets[idx : idx + MAX_PPTX_BULLETS_PER_SLIDE]),
                }
            )
    return expanded[:MAX_PPTX_SLIDES]


def _set_shape_text(shape: Any, text: str, font_size: int = 20, bold: bool = False) -> None:
    try:
        from pptx.dml.color import RGBColor as PptxRGBColor  # type: ignore

        brand_blue = _brand_value("colors", "blue", "0069B3")
        text_frame = shape.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text or ""
        run.font.size = Pt(font_size)
        run.font.bold = bool(bold)
        if bold:
            run.font.color.rgb = PptxRGBColor(
                int(brand_blue[0:2], 16), int(brand_blue[2:4], 16), int(brand_blue[4:6], 16)
            )
    except Exception:
        try:
            shape.text = text or ""
        except Exception:
            pass


def _shape_key(shape: Any) -> int:
    try:
        return int(shape.shape_id)
    except Exception:
        return id(shape)


def _shape_area(shape: Any) -> int:
    try:
        return int(shape.width) * int(shape.height)
    except Exception:
        return 0


def _clear_shape_text(shape: Any) -> None:
    try:
        shape.text_frame.clear()
    except Exception:
        try:
            shape.text = ""
        except Exception:
            pass


def _prepare_template_slide_text_shapes(slide: Any) -> tuple[Any, Any]:
    text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]

    title_shape = None
    try:
        if slide.shapes.title is not None and getattr(slide.shapes.title, "has_text_frame", False):
            title_shape = slide.shapes.title
    except Exception:
        title_shape = None

    if title_shape is None and text_shapes:
        title_shape = sorted(
            text_shapes,
            key=lambda shape: (getattr(shape, "top", 0), -_shape_area(shape)),
        )[0]

    title_key = _shape_key(title_shape) if title_shape is not None else -1
    body_candidates = [shape for shape in text_shapes if _shape_key(shape) != title_key]
    body_shape = max(body_candidates, key=_shape_area) if body_candidates else None

    keep = {_shape_key(shape) for shape in (title_shape, body_shape) if shape is not None}
    for shape in text_shapes:
        if _shape_key(shape) not in keep:
            _clear_shape_text(shape)

    return title_shape, body_shape


PPTX_PLACEHOLDER_TEXTS = (
    "Click to edit Master title style",
    "Click to edit Master subtitle style",
    "Click to edit Master text styles",
    "Mastertextformat bearbeiten",
    "Zweite Ebene",
    "Dritte Ebene",
    "Vierte Ebene",
    "Fuenfte Ebene",
    "Fünfte Ebene",
)


def _sanitize_pptx_template_artifacts(data: bytes) -> bytes:
    """
    Some PowerPoint templates contain visible sample text in slide masters.
    Generated slides inherit that text unless it is removed from the package.
    """
    try:
        src = io.BytesIO(data)
        out = io.BytesIO()
        with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                payload = zin.read(item.filename)
                if item.filename.startswith(("ppt/slideMasters/", "ppt/slideLayouts/")) and item.filename.endswith(".xml"):
                    for placeholder in PPTX_PLACEHOLDER_TEXTS:
                        payload = payload.replace(
                            f"<a:t>{html.escape(placeholder)}</a:t>".encode("utf-8"),
                            b"<a:t></a:t>",
                        )
                        payload = payload.replace(
                            f"<a:t>{placeholder}</a:t>".encode("utf-8"),
                            b"<a:t></a:t>",
                        )
                zout.writestr(item, payload)
        return out.getvalue()
    except Exception:
        return data


def _markdown_to_pptx_bytes(content: str, title: str = "Praesentation") -> Optional[bytes]:
    if Presentation is None or Inches is None or Pt is None:
        return None

    try:
        if _file_exists(KAHLE_PPTX_TEMPLATE):
            prs = Presentation(str(KAHLE_PPTX_TEMPLATE))
            _PPTX_TEMPLATE_USED.set(True)
        else:
            prs = Presentation()

        prs.core_properties.author = "KAHLE-Vinci"
        prs.core_properties.title = (title or "Praesentation").strip() or "Praesentation"

        specs = _markdown_to_slide_specs(content, title)
        for index, spec in enumerate(specs):
            if index < len(prs.slides):
                slide = prs.slides[index]
            else:
                layout = _pick_pptx_layout(prs, ("title", "titel") if index == 0 else ("content", "inhalt", "title"))
                slide = prs.slides.add_slide(layout)

            title_shape, body_shape = _prepare_template_slide_text_shapes(slide)
            if title_shape is not None:
                _set_shape_text(title_shape, spec["title"], font_size=26 if index == 0 else 22, bold=True)
            else:
                title_shape = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(8.4), Inches(0.7))
                _set_shape_text(title_shape, spec["title"], font_size=26 if index == 0 else 22, bold=True)

            if body_shape is None:
                body_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(8.2), Inches(4.6))

            text_frame = body_shape.text_frame
            text_frame.clear()
            bullets = spec.get("bullets", []) or [" "]
            for bullet_index, bullet in enumerate(bullets):
                paragraph = text_frame.paragraphs[0] if bullet_index == 0 else text_frame.add_paragraph()
                paragraph.text = str(bullet)
                paragraph.level = 0
                for run in paragraph.runs:
                    run.font.size = Pt(18 if index == 0 else 15)

            if index > 0:
                try:
                    footer = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(8.5), Inches(0.25))
                    footer.text = "KAHLE-Vinci"
                    footer.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
                except Exception:
                    pass

        _remove_pptx_slides_after(prs, len(specs))

        out = io.BytesIO()
        prs.save(out)
        return _sanitize_pptx_template_artifacts(out.getvalue())
    except Exception:
        return None


def _sign_download(rel: str, exp: int) -> str:
    if not FILE_LINK_SECRET:
        raise HTTPException(status_code=500, detail="FILE_LINK_SECRET not set")
    msg = f"{exp}:{rel}".encode("utf-8")
    return hmac.new(FILE_LINK_SECRET.encode("utf-8"), msg, hashlib.sha256).hexdigest()


def _verify_sig(rel: str, exp: int, sig: str) -> None:
    if exp < int(time.time()):
        raise HTTPException(status_code=410, detail="link_expired")
    if ALLOW_UNSIGNED_DOWNLOADS:
        return
    expected = _sign_download(rel, exp)
    if not hmac.compare_digest(expected, sig):
        raise HTTPException(status_code=401, detail="bad_signature")


def _encode_download_token(rel: str, exp: int, sig: str) -> str:
    payload = json.dumps({"rel": rel, "exp": exp, "sig": sig}, separators=(",", ":")).encode("utf-8")
    return base64.urlsafe_b64encode(payload).decode("ascii").rstrip("=")


def _decode_download_token(token: str) -> tuple[str, int, str]:
    try:
        padded = token + ("=" * (-len(token) % 4))
        raw = base64.urlsafe_b64decode(padded.encode("ascii"))
        payload = json.loads(raw.decode("utf-8"))
        rel = _safe_relpath(str(payload.get("rel") or ""))
        exp = int(payload.get("exp"))
        sig = str(payload.get("sig") or "")
    except Exception:
        raise HTTPException(status_code=401, detail="bad_download_token")
    if not rel or not sig:
        raise HTTPException(status_code=401, detail="bad_download_token")
    return rel, exp, sig


def _build_download_url(rel: str) -> str:
    exp = int(time.time()) + FILE_LINK_TTL_SECONDS
    sig = "unsigned" if ALLOW_UNSIGNED_DOWNLOADS else _sign_download(rel, exp)
    token = _encode_download_token(rel, exp, sig)
    path = f"/files/download?token={token}"
    return f"{PUBLIC_BASE_URL}{path}" if PUBLIC_BASE_URL else path


def _save_bytes(filename: str, data: bytes) -> dict:
    safe_name = _sanitize_filename(filename)
    ext = Path(safe_name).suffix.lower()
    if ext and ext not in ALLOWED_SAVE_EXT:
        raise HTTPException(status_code=400, detail=f"extension_not_allowed: {ext}")

    ts = time.strftime("%Y%m%d_%H%M%S")
    rand = uuid.uuid4().hex[:8]
    rel = _safe_relpath(f"{OUTPUT_SUBDIR}/{ts}_{rand}_{safe_name}")

    abs_path = (UPLOAD_ROOT / rel).resolve()
    _ensure_within_upload_root(abs_path)

    _write_atomic(abs_path, data)

    return {
        "output_kind": "file_saved",
        "saved_rel_path": rel,
        "filename": safe_name,
        "size_bytes": len(data),
        "sha256": _sha256(data),
        "download_url": _build_download_url(rel),
    }


# -----------------------------
# Request models
# -----------------------------
class SaveB64Request(BaseModel):
    filename: str = Field(..., description="Output filename, e.g. edited.docx")
    content_type: str = Field(..., description="MIME type")
    content_base64: str = Field(..., description="Base64 content (no data: prefix)")


class ReplaceOneSaveRequest(BaseModel):
    from_text: str = Field(..., description="Text to find (first occurrence)")
    to_text: str = Field(..., description="Replacement text")
    file_path: NonEmptyStr = Field(..., description="Exact filename/path in uploads")


class BundleToMdRequest(BaseModel):
    title: str = Field("Masterkontext", description="Title for the bundle output")
    file_paths: Optional[list[str]] = Field(None, description="List of exact filenames/paths in uploads")

class BundleToMdSaveRequest(BaseModel):
    title: str = Field("Masterkontext", description="Title for the bundle output")
    file_paths: list[NonEmptyStr] = Field(..., min_length=1, description="List of exact filenames/paths in uploads")


class DocxDeleteLastParagraphsSaveRequest(BaseModel):
    file_path: NonEmptyStr = Field(..., description="Exact DOCX filename in uploads")
    n: int = Field(3, ge=1, le=500, description="How many paragraphs to delete from the end")
    allow_empty_output: bool = Field(
        False,
        description="If false, block operations that would remove all non-empty DOCX paragraphs.",
    )


class TextApplyOpsSaveRequest(BaseModel):
    file_path: NonEmptyStr = Field(..., description="Exact TXT/MD/CSV filename in uploads")
    ops: list[dict[str, Any]] = Field(..., description="List of deterministic text ops")
    allow_empty_output: bool = Field(
        False,
        description="If false, block operations that would result in a 0-byte file.",
    )


class PdfRemovePagesSaveRequest(BaseModel):
    file_path: NonEmptyStr = Field(..., description="Exact PDF filename in uploads")
    remove_pages: Optional[list[int]] = Field(None, description="1-based page numbers to remove, e.g. [1,2]")
    pages_to_remove: Optional[list[int]] = Field(None, description="Alias for remove_pages.")
    remove_last_page: bool = Field(False, description="If true, remove the final page of the PDF.")


class PdfMergeSaveRequest(BaseModel):
    file_paths: list[NonEmptyStr] = Field(..., min_length=2, description="At least 2 exact PDF filenames in uploads")
    output_name: str = Field("merged.pdf", description="Output filename")
    attachment_exact_paths: Optional[list[str]] = Field(
        default=None,
        description="Internal: exact attachment paths from current OWUI message for disambiguation.",
    )


class XlsxUpdateCellsSaveRequest(BaseModel):
    file_path: NonEmptyStr = Field(..., description="Exact XLSX filename in uploads")
    updates: list[dict[str, Any]] = Field(
        ...,
        description='Cell/range updates, e.g. [{"sheet":"Sheet1","cell":"B2","value":"123"}] or [{"range":"A2:A500","generator":"random_money","min":1000,"max":100000}]',
    )


class DocxToPdfSaveRequest(BaseModel):
    file_path: NonEmptyStr = Field(..., description="Exact DOCX filename in uploads")
    output_name: str = Field("converted.pdf", description="Output PDF filename")


class FileToMdSaveRequest(BaseModel):
    file_path: NonEmptyStr = Field(..., description="Exact filename in uploads (docx/pdf/md/txt/xlsx/csv)")
    title: str = Field("Konvertiert", description="Heading title for markdown conversion")
    output_name: Optional[NonEmptyStr] = Field(None, description="Optional output filename, defaults to source stem + .md")


class FileToDocxSaveRequest(BaseModel):
    file_path: NonEmptyStr = Field(..., description="Exact filename in uploads (docx/pdf/md/txt/xlsx/csv)")
    title: str = Field("Konvertiert", description="Heading title for DOCX conversion")
    output_name: Optional[NonEmptyStr] = Field(None, description="Optional output filename, defaults to source stem + .docx")


class TextCreateSaveRequest(BaseModel):
    filename: NonEmptyStr = Field(..., description="Output filename ending in .md, .txt, or .csv")
    content: str = Field(..., description="Text/Markdown/CSV content to save")


class DocxCreateSaveRequest(BaseModel):
    filename: NonEmptyStr = Field(..., description="Output filename ending in .docx")
    content: str = Field(
        ...,
        description=(
            "Markdown/text content to render into a simple DOCX. "
            "When the user asks to create a document from the result, pass the full previous assistant result here. "
            "Never leave this empty."
        ),
    )
    title: str = Field("Dokument", description="Document title")


class PdfCreateSaveRequest(BaseModel):
    filename: NonEmptyStr = Field(..., description="Output filename ending in .pdf")
    content: str = Field(
        ...,
        description=(
            "Markdown/text content to render into a simple PDF. "
            "When the user asks to create a PDF from the result, pass the full previous assistant result here. "
            "Never leave this empty."
        ),
    )
    title: str = Field("Dokument", description="Document title")


class PptxCreateSaveRequest(BaseModel):
    filename: NonEmptyStr = Field(..., description="Output filename ending in .pptx")
    content: str = Field(
        ...,
        description=(
            "Markdown/text content to render into a PowerPoint deck. "
            "Use this when the user asks to create a presentation or PPTX from generated content. "
            "Never leave this empty."
        ),
    )
    title: str = Field("Praesentation", description="Presentation title")


class CleanupOldFilesRequest(BaseModel):
    days: int = Field(15, ge=1, le=3650, description="Delete files older than this many days")
    dry_run: bool = Field(False, description="If true, list files only, do not delete")
    include_uploads: bool = Field(True, description="Include /uploads root")
    include_edited: bool = Field(True, description="Include /edited output folder")
    include_legacy_edited: bool = Field(
        True,
        description="Also include legacy /app/backend/data/edited when uploads root is /app/backend/data/uploads",
    )
    max_files: int = Field(100000, ge=1, le=1000000, description="Safety limit for processed files")


class CleanupOpenWebUIChatsRequest(BaseModel):
    days: int = Field(60, ge=1, le=3650, description="Delete chats older than this many days")
    dry_run: bool = Field(False, description="If true, only return candidates")
    keep_pinned: bool = Field(True, description="If true, never delete pinned chats")
    max_delete: int = Field(5000, ge=1, le=100000, description="Safety limit for deleted chats per run")
    vacuum: bool = Field(True, description="If true, run VACUUM after deletion")


# -----------------------------
# Endpoints
# -----------------------------
@app.get("/health", include_in_schema=False)
def health():
    return {
        "ok": True,
        "owui_upload_dir": str(UPLOAD_ROOT),
        "doc_worker_url": DOC_WORKER_URL,
        "output_subdir": OUTPUT_SUBDIR,
        "max_file_bytes": MAX_FILE_BYTES,
        "signed_downloads_enabled": bool(FILE_LINK_SECRET) and not ALLOW_UNSIGNED_DOWNLOADS,
        "allow_unsigned_downloads": ALLOW_UNSIGNED_DOWNLOADS,
        "public_base_url": PUBLIC_BASE_URL or None,
        "tool_api_key_configured": bool(TOOL_API_KEY),
        "require_tool_api_key": REQUIRE_TOOL_API_KEY,
        "require_exact_file_path": REQUIRE_EXACT_FILE_PATH,
        "disallow_wildcards": DISALLOW_WILDCARDS,
        "recent_upload_disambiguation_seconds": RECENT_UPLOAD_DISAMBIGUATION_SECONDS,
        "assets_root": str(KAHLE_ASSETS_ROOT),
        "templates": {
            "docx": _file_exists(KAHLE_DOCX_TEMPLATE),
            "pdf": _file_exists(KAHLE_PDF_TEMPLATE),
            "pptx": _file_exists(KAHLE_PPTX_TEMPLATE),
            "logo": _file_exists(KAHLE_LOGO_PRIMARY),
        },
        "renderers": {
            "python_docx": Document is not None,
            "python_pptx": Presentation is not None,
            "reportlab": SimpleDocTemplate is not None,
            "pypdf": PdfReader is not None and PdfWriter is not None,
        },
    }


def _cleanup_dir_old_files(target: Path, cutoff_ts: float, dry_run: bool, max_files: int) -> dict[str, Any]:
    result: dict[str, Any] = {
        "path": str(target),
        "exists": target.exists() and target.is_dir(),
        "checked_files": 0,
        "matched_files": 0,
        "deleted_files": 0,
        "deleted_bytes": 0,
        "samples": [],
        "truncated": False,
    }
    if not result["exists"]:
        return result

    processed = 0
    for p in target.rglob("*"):
        if not p.is_file():
            continue
        processed += 1
        if processed > max_files:
            result["truncated"] = True
            break
        result["checked_files"] += 1
        try:
            st = p.stat()
        except Exception:
            continue
        if float(st.st_mtime) >= cutoff_ts:
            continue

        result["matched_files"] += 1
        try:
            rel = str(p.resolve().relative_to(UPLOAD_ROOT)).replace("\\", "/")
        except Exception:
            rel = str(p.name)
        if len(result["samples"]) < 200:
            result["samples"].append(rel)

        if dry_run:
            continue
        try:
            size = int(st.st_size)
            p.unlink(missing_ok=True)
            result["deleted_files"] += 1
            result["deleted_bytes"] += size
        except Exception:
            continue

    if not dry_run:
        # Clean empty subdirectories but keep the directory itself.
        for d in sorted([x for x in target.rglob("*") if x.is_dir()], key=lambda x: len(x.parts), reverse=True):
            try:
                if not any(d.iterdir()):
                    d.rmdir()
            except Exception:
                continue

    return result


def _parse_iso_ts(value: Any) -> Optional[datetime]:
    if not isinstance(value, str) or not value.strip():
        return None
    raw = value.strip()
    if raw.endswith("Z"):
        raw = raw[:-1] + "+00:00"
    try:
        dt = datetime.fromisoformat(raw)
    except Exception:
        return None
    if dt.tzinfo is None:
        dt = dt.replace(tzinfo=timezone.utc)
    return dt.astimezone(timezone.utc)


def _load_log_cleanup_report() -> dict[str, Any]:
    if not OWUI_LOG_CLEANUP_REPORT_PATH.exists() or not OWUI_LOG_CLEANUP_REPORT_PATH.is_file():
        return {
            "ok": True,
            "exists": False,
            "report_path": str(OWUI_LOG_CLEANUP_REPORT_PATH),
            "success": False,
            "stale": True,
            "reason": "report_not_found",
        }
    try:
        payload = json.loads(OWUI_LOG_CLEANUP_REPORT_PATH.read_text(encoding="utf-8-sig"))
    except Exception as e:
        return {
            "ok": True,
            "exists": True,
            "report_path": str(OWUI_LOG_CLEANUP_REPORT_PATH),
            "success": False,
            "stale": True,
            "reason": "invalid_report_json",
            "error": str(e),
        }

    last_run = payload.get("last_run_at")
    run_dt = _parse_iso_ts(last_run)
    now_dt = datetime.now(timezone.utc)
    stale = True
    age_hours: Optional[float] = None
    if run_dt:
        age_hours = max(0.0, (now_dt - run_dt).total_seconds() / 3600.0)
        stale = age_hours > float(OWUI_LOG_CLEANUP_STALE_HOURS)

    return {
        "ok": True,
        "exists": True,
        "report_path": str(OWUI_LOG_CLEANUP_REPORT_PATH),
        "last_run_at": last_run,
        "success": bool(payload.get("success", False)),
        "cutoff_days": int(payload.get("cutoff_days", 0) or 0),
        "lines_before": int(payload.get("lines_before", 0) or 0),
        "lines_after": int(payload.get("lines_after", 0) or 0),
        "lines_deleted": int(payload.get("lines_deleted", 0) or 0),
        "open_webui_restarted": bool(payload.get("open_webui_restarted", False)),
        "stale": stale,
        "age_hours": age_hours,
        "error": payload.get("error"),
    }


@app.post("/maintenance/cleanup_old_files", include_in_schema=False)
def maintenance_cleanup_old_files(
    payload: CleanupOldFilesRequest,
    x_api_key: Optional[str] = Header(default=None),
):
    _require_api_key(x_api_key)

    now_ts = time.time()
    cutoff_ts = now_ts - (payload.days * 86400)

    targets: list[Path] = []
    if payload.include_uploads:
        targets.append(UPLOAD_ROOT)
    # /edited lives under /uploads, so avoid double-processing when uploads is selected.
    if payload.include_edited and not payload.include_uploads:
        edited_dir = (UPLOAD_ROOT / OUTPUT_SUBDIR).resolve()
        _ensure_within_upload_root(edited_dir)
        if edited_dir not in targets:
            targets.append(edited_dir)
    if payload.include_legacy_edited:
        base_parent = UPLOAD_ROOT.parent.resolve()
        legacy_edited_dir = (base_parent / OUTPUT_SUBDIR).resolve()
        try:
            legacy_edited_dir.relative_to(base_parent)
        except Exception:
            legacy_edited_dir = None
        if (
            legacy_edited_dir
            and legacy_edited_dir != UPLOAD_ROOT
            and legacy_edited_dir not in targets
        ):
            targets.append(legacy_edited_dir)

    reports = []
    for t in targets:
        reports.append(
            _cleanup_dir_old_files(
                target=t,
                cutoff_ts=cutoff_ts,
                dry_run=payload.dry_run,
                max_files=payload.max_files,
            )
        )

    summary = {
        "ok": True,
        "dry_run": payload.dry_run,
        "days": payload.days,
        "cutoff_unix": int(cutoff_ts),
        "targets_count": len(reports),
        "checked_files": sum(int(r.get("checked_files", 0)) for r in reports),
        "matched_files": sum(int(r.get("matched_files", 0)) for r in reports),
        "deleted_files": sum(int(r.get("deleted_files", 0)) for r in reports),
        "deleted_bytes": sum(int(r.get("deleted_bytes", 0)) for r in reports),
        "reports": reports,
    }
    return summary


@app.post("/maintenance/cleanup_openwebui_chats", include_in_schema=False)
def maintenance_cleanup_openwebui_chats(
    payload: CleanupOpenWebUIChatsRequest,
    x_api_key: Optional[str] = Header(default=None),
):
    _require_api_key(x_api_key)
    cutoff_ts = int(time.time() - (payload.days * 86400))

    if not OWUI_DB_PATH.exists() or not OWUI_DB_PATH.is_file():
        raise HTTPException(status_code=500, detail=f"openwebui_db_not_found: {OWUI_DB_PATH}")

    conn = sqlite3.connect(str(OWUI_DB_PATH), timeout=30.0)
    try:
        conn.row_factory = sqlite3.Row

        where_sql = "CAST(updated_at AS INTEGER) < ?"
        where_args: list[Any] = [cutoff_ts]
        if payload.keep_pinned:
            where_sql += " AND COALESCE(pinned, 0) = 0"

        candidates = int(conn.execute(f"SELECT COUNT(*) AS c FROM chat WHERE {where_sql}", where_args).fetchone()["c"])
        skipped_pinned = 0
        if payload.keep_pinned:
            skipped_pinned = int(
                conn.execute(
                    "SELECT COUNT(*) AS c FROM chat WHERE CAST(updated_at AS INTEGER) < ? AND COALESCE(pinned, 0) != 0",
                    (cutoff_ts,),
                ).fetchone()["c"]
            )

        rows = conn.execute(
            f"SELECT id FROM chat WHERE {where_sql} ORDER BY CAST(updated_at AS INTEGER) ASC LIMIT ?",
            [*where_args, payload.max_delete],
        ).fetchall()
        candidate_ids = [str(r["id"]) for r in rows if r["id"]]

        deleted_tags = 0
        deleted_chats = 0
        vacuum_ran = False

        if candidate_ids and not payload.dry_run:
            placeholders = ",".join("?" for _ in candidate_ids)
            conn.execute("BEGIN")
            try:
                deleted_tags = int(
                    conn.execute(
                        f"DELETE FROM chatidtag WHERE chat_id IN ({placeholders})",
                        candidate_ids,
                    ).rowcount
                    or 0
                )
                deleted_chats = int(
                    conn.execute(
                        f"DELETE FROM chat WHERE id IN ({placeholders})",
                        candidate_ids,
                    ).rowcount
                    or 0
                )
                conn.commit()
            except Exception:
                conn.rollback()
                raise

            if payload.vacuum and deleted_chats > 0:
                conn.execute("VACUUM")
                vacuum_ran = True

        return {
            "ok": True,
            "dry_run": payload.dry_run,
            "days": payload.days,
            "keep_pinned": payload.keep_pinned,
            "max_delete": payload.max_delete,
            "cutoff_unix": cutoff_ts,
            "candidates": candidates,
            "selected_for_delete": len(candidate_ids),
            "deleted_chats": deleted_chats,
            "deleted_tags": deleted_tags,
            "skipped_pinned": skipped_pinned,
            "sample_chat_ids": candidate_ids[:20],
            "vacuum_ran": vacuum_ran,
        }
    finally:
        conn.close()


@app.get("/maintenance/openwebui_log_cleanup_status", include_in_schema=False)
def maintenance_openwebui_log_cleanup_status(x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)
    return _load_log_cleanup_report()


@app.post("/files/save_b64", include_in_schema=False)
def files_save_b64(payload: SaveB64Request, x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)
    data = _b64decode_strict(payload.content_base64)
    saved = _save_bytes(payload.filename, data)
    saved["content_type"] = payload.content_type
    return saved


@app.post("/text/create_save", operation_id="text_create_save", include_in_schema=False)
def text_create_save(payload: TextCreateSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    """
    Create a downloadable Markdown/TXT/CSV file from model-generated text.
    Intended for research briefs, summaries, and structured Markdown outputs.
    """
    _require_api_key(x_api_key)

    safe_name = _sanitize_filename(payload.filename)
    ext = Path(safe_name).suffix.lower()
    if ext not in {".md", ".txt", ".csv"}:
        raise HTTPException(status_code=400, detail="text_create_save_allows_only_md_txt_csv")

    data = payload.content.encode("utf-8")
    if len(data) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail=f"file_too_large (>{MAX_FILE_BYTES} bytes)")

    saved = _save_bytes(safe_name, data)
    saved["content_type"] = "text/markdown; charset=utf-8" if ext == ".md" else "text/plain; charset=utf-8"
    return saved


@app.post("/docx/create_save", operation_id="docx_create_save", include_in_schema=False)
def docx_create_save(payload: DocxCreateSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    """
    Create a simple downloadable DOCX from model-generated Markdown/text.
    Use this after research when the user asks for a DOCX file.
    """
    _require_api_key(x_api_key)

    out_name = _sanitize_filename(payload.filename)
    if not out_name.lower().endswith(".docx"):
        out_name = f"{out_name}.docx"

    _DOCX_TEMPLATE_USED.set(False)
    data = _markdown_to_docx_bytes(payload.content, payload.title)
    if len(data) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail=f"file_too_large (>{MAX_FILE_BYTES} bytes)")

    saved = _save_bytes(out_name, data)
    saved["content_type"] = DOCX_MIME
    saved["conversion"] = "markdown_to_docx"
    saved["template_used"] = bool(_DOCX_TEMPLATE_USED.get())
    return saved


@app.post("/pdf/create_save", operation_id="pdf_create_save", include_in_schema=False)
def pdf_create_save(payload: PdfCreateSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    """
    Create a simple downloadable PDF from model-generated Markdown/text.
    Use this after research when the user asks for a PDF file.
    """
    _require_api_key(x_api_key)

    out_name = _sanitize_filename(payload.filename)
    if not out_name.lower().endswith(".pdf"):
        out_name = f"{out_name}.pdf"

    _PDF_TEMPLATE_USED.set(False)
    data = _text_to_pdf_bytes(payload.content, payload.title)
    if len(data) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail=f"file_too_large (>{MAX_FILE_BYTES} bytes)")

    saved = _save_bytes(out_name, data)
    saved["content_type"] = PDF_MIME
    saved["conversion"] = "markdown_to_pdf"
    saved["template_used"] = bool(_PDF_TEMPLATE_USED.get())
    return saved


@app.post("/pptx/create_save", operation_id="pptx_create_save", include_in_schema=False)
def pptx_create_save(payload: PptxCreateSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    """
    Create a downloadable PPTX deck from model-generated Markdown/text.
    Uses the KAHLE PowerPoint template when available.
    """
    _require_api_key(x_api_key)

    out_name = _sanitize_filename(payload.filename)
    if not out_name.lower().endswith(".pptx"):
        out_name = f"{out_name}.pptx"

    _PPTX_TEMPLATE_USED.set(False)
    data = _markdown_to_pptx_bytes(payload.content, payload.title)
    if not data:
        raise HTTPException(status_code=500, detail="pptx_renderer_unavailable")
    if len(data) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail=f"file_too_large (>{MAX_FILE_BYTES} bytes)")

    saved = _save_bytes(out_name, data)
    saved["content_type"] = PPTX_MIME
    saved["conversion"] = "markdown_to_pptx"
    saved["template_used"] = bool(_PPTX_TEMPLATE_USED.get())
    return saved


@app.get("/files/download", include_in_schema=False)
def files_download(
    rel: Optional[str] = Query(default=None, description="Relative path under uploads"),
    exp: Optional[int] = Query(default=None, description="Expiry (unix timestamp)"),
    sig: Optional[str] = Query(default=None, description="HMAC signature"),
    token: Optional[str] = Query(default=None, description="Opaque signed download token"),
    x_api_key: Optional[str] = Header(default=None),
    authorization: Optional[str] = Header(default=None),
):
    if token:
        rel, exp, sig = _decode_download_token(token)
    elif rel:
        rel = _safe_relpath(rel)
    else:
        raise HTTPException(status_code=400, detail="download_token_or_rel_required")

    if not ALLOW_UNSIGNED_DOWNLOADS:
        # Compatibility fallback:
        # Some internal tool callers may accidentally keep only `rel`. Only authenticated
        # callers may mint a fresh short-lived signed link.
        if exp is None or not sig:
            _require_api_key(x_api_key, authorization)
            signed_url = _build_download_url(rel)
            return RedirectResponse(url=signed_url, status_code=307)

        _verify_sig(rel, exp, sig)

    abs_path = (UPLOAD_ROOT / rel).resolve()
    _ensure_within_upload_root(abs_path)

    if not abs_path.exists() or not abs_path.is_file():
        raise HTTPException(status_code=404, detail="file_not_found")

    media_type, _ = mimetypes.guess_type(abs_path.name)
    media_type = media_type or "application/octet-stream"

    return FileResponse(path=str(abs_path), filename=abs_path.name, media_type=media_type)


# -----------------------------
# DOCX: replace_one (existing)
# -----------------------------
@app.post("/docx/replace_one_b64", include_in_schema=False)
def docx_replace_one_b64(
    from_text: str,
    to_text: str,
    file_path: str | None = None,
    x_api_key: Optional[str] = Header(default=None),
):
    _require_api_key(x_api_key)

    path = _resolve_path(file_path, allowed_exts=(".docx",))
    data = _read_file(path)

    files = {"file": (path.name, data, DOCX_MIME)}
    form = {"from_text": from_text, "to_text": to_text}

    r = requests.post(
        f"{DOC_WORKER_URL}/docx/replace_one",
        headers=_worker_headers(),
        files=files,
        data=form,
        timeout=120,
    )
    _raise_for_worker_response(r)

    b64 = base64.b64encode(r.content).decode("ascii")
    out_name = f"edited_{path.name}"
    return JSONResponse(content={"filename": out_name, "content_type": DOCX_MIME, "content_base64": b64})


@app.post("/docx/replace_one_save", operation_id="docx_replace_one_save")
def docx_replace_one_save(payload: ReplaceOneSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)

    path = _resolve_path(payload.file_path, allowed_exts=(".docx",), require_exact=True)
    data = _read_file(path)

    files = {"file": (path.name, data, DOCX_MIME)}
    form = {"from_text": payload.from_text, "to_text": payload.to_text}

    r = requests.post(
        f"{DOC_WORKER_URL}/docx/replace_one",
        headers=_worker_headers(),
        files=files,
        data=form,
        timeout=180,
    )
    _raise_for_worker_response(r)

    out_name = f"edited_{path.name}"
    saved = _save_bytes(out_name, r.content)
    saved["content_type"] = DOCX_MIME
    saved["source_file"] = str(path.name)
    return saved


@app.post("/docx/delete_last_paragraphs_save", operation_id="docx_delete_last_paragraphs_save")
def docx_delete_last_paragraphs_save(
    payload: DocxDeleteLastParagraphsSaveRequest,
    x_api_key: Optional[str] = Header(default=None),
):
    _require_api_key(x_api_key)

    path = _resolve_path(payload.file_path, allowed_exts=(".docx",), require_exact=True)
    data = _read_file(path)

    files = {"file": (path.name, data, DOCX_MIME)}
    form = {
        "n": str(payload.n),
        "allow_empty_output": "true" if payload.allow_empty_output else "false",
    }

    r = requests.post(
        f"{DOC_WORKER_URL}/docx/delete_last_paragraphs",
        headers=_worker_headers(),
        files=files,
        data=form,
        timeout=180,
    )
    _raise_for_worker_response(r)

    out_name = f"edited_{path.name}"
    saved = _save_bytes(out_name, r.content)
    saved["content_type"] = DOCX_MIME
    saved["source_file"] = str(path.name)
    saved["n_deleted_paragraphs"] = payload.n
    return saved


# -----------------------------
# TEXT: apply_ops_save (new)
# -----------------------------
@app.post("/text/apply_ops_save", operation_id="text_apply_ops_save")
def text_apply_ops_save(payload: TextApplyOpsSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)

    path = _resolve_path(payload.file_path, allowed_exts=(".txt", ".md", ".csv"), require_exact=True)
    data = _read_file(path)

    ext = path.suffix.lower()
    if ext not in (".txt", ".md", ".csv"):
        raise HTTPException(status_code=400, detail="invalid_text_extension")

    files = {"file": (path.name, data, "application/octet-stream")}
    form = {"ops_json": json.dumps(payload.ops, ensure_ascii=False)}

    r = requests.post(
        f"{DOC_WORKER_URL}/text/apply_ops",
        headers=_worker_headers(),
        files=files,
        data=form,
        timeout=180,
    )
    _raise_for_worker_response(r)
    if (not payload.allow_empty_output) and len(r.content) == 0:
        raise HTTPException(
            status_code=400,
            detail="empty_output_blocked_set_allow_empty_output_true_to_override",
        )

    out_name = f"edited_{path.name}"
    saved = _save_bytes(out_name, r.content)
    saved["content_type"] = "text/markdown; charset=utf-8" if ext == ".md" else "text/plain; charset=utf-8"
    saved["source_file"] = str(path.name)
    saved["ops_count"] = len(payload.ops)
    return saved


# -----------------------------
# XLSX: update_cells_save
# -----------------------------
@app.post("/xlsx/update_cells_save", operation_id="xlsx_update_cells_save")
def xlsx_update_cells_save(payload: XlsxUpdateCellsSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)

    if not payload.updates:
        raise HTTPException(status_code=400, detail="updates_required_non_empty")

    path = _resolve_path(payload.file_path, allowed_exts=(".xlsx",), require_exact=True)
    data = _read_file(path)

    files = {"file": (path.name, data, XLSX_MIME)}
    form = {"updates_json": json.dumps(payload.updates, ensure_ascii=False)}

    r = requests.post(
        f"{DOC_WORKER_URL}/xlsx/update_cells",
        headers=_worker_headers(),
        files=files,
        data=form,
        timeout=180,
    )
    _raise_for_worker_response(r)

    out_name = f"edited_{path.name}"
    saved = _save_bytes(out_name, r.content)
    saved["content_type"] = XLSX_MIME
    saved["source_file"] = str(path.name)
    saved["updates_count"] = len(payload.updates)
    return saved


# -----------------------------
# DOCX: to_pdf_save
# -----------------------------
@app.post("/docx/to_pdf_save", operation_id="docx_to_pdf_save")
def docx_to_pdf_save(payload: DocxToPdfSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)

    path = _resolve_path(payload.file_path, allowed_exts=(".docx",), require_exact=True)
    data = _read_file(path)

    out_name = payload.output_name or f"{Path(path.name).stem}.pdf"
    if not out_name.lower().endswith(".pdf"):
        out_name = f"{out_name}.pdf"

    title = Path(path.name).stem or "Konvertiert"
    files = [("files", (path.name, data, DOCX_MIME))]
    r = requests.post(
        f"{DOC_WORKER_URL}/bundle/to_md",
        headers=_worker_headers(),
        files=files,
        data={"title": title, "mode": "raw"},
        timeout=300,
    )
    _raise_for_worker_response(r)

    markdown = r.content.decode("utf-8", errors="replace")
    _PDF_TEMPLATE_USED.set(False)
    out = _text_to_pdf_bytes(markdown, title)
    if len(out) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail=f"file_too_large (>{MAX_FILE_BYTES} bytes)")

    saved = _save_bytes(out_name, out)
    saved["content_type"] = PDF_MIME
    saved["source_file"] = str(path.name)
    saved["conversion"] = "docx_to_pdf_branded"
    saved["template_used"] = bool(_PDF_TEMPLATE_USED.get())
    return saved


# -----------------------------
# Generic single-file -> Markdown
# -----------------------------
@app.post("/file/to_md_save", operation_id="file_to_md_save")
def file_to_md_save(payload: FileToMdSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)

    path = _resolve_path(
        payload.file_path,
        allowed_exts=(".docx", ".pdf", ".md", ".txt", ".xlsx", ".csv"),
        require_exact=True,
    )
    data = _read_file(path)

    title = (payload.title or "Konvertiert").strip() or "Konvertiert"
    mfiles = [("files", (path.name, data, "application/octet-stream"))]

    r = requests.post(
        f"{DOC_WORKER_URL}/bundle/to_md",
        headers=_worker_headers(),
        files=mfiles,
        data={"title": title, "mode": "raw"},
        timeout=300,
    )
    _raise_for_worker_response(r)

    if payload.output_name:
        out_name = payload.output_name
    else:
        out_name = f"{Path(path.name).stem}.md"
    if not out_name.lower().endswith(".md"):
        out_name = f"{out_name}.md"

    saved = _save_bytes(out_name, r.content)
    saved["content_type"] = "text/markdown; charset=utf-8"
    saved["source_file"] = str(path.name)
    saved["conversion"] = "file_to_md"
    return saved


@app.post("/file/to_docx_save", operation_id="file_to_docx_save")
def file_to_docx_save(payload: FileToDocxSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)

    path = _resolve_path(
        payload.file_path,
        allowed_exts=(".docx", ".pdf", ".md", ".txt", ".xlsx", ".csv"),
        require_exact=True,
    )
    data = _read_file(path)

    title = (payload.title or "Konvertiert").strip() or "Konvertiert"
    mfiles = [("files", (path.name, data, "application/octet-stream"))]

    r = requests.post(
        f"{DOC_WORKER_URL}/bundle/to_md",
        headers=_worker_headers(),
        files=mfiles,
        data={"title": title, "mode": "raw"},
        timeout=300,
    )
    _raise_for_worker_response(r)

    if payload.output_name:
        out_name = payload.output_name
    else:
        out_name = f"{Path(path.name).stem}.docx"
    if not out_name.lower().endswith(".docx"):
        out_name = f"{out_name}.docx"

    markdown = r.content.decode("utf-8", errors="replace")
    _DOCX_TEMPLATE_USED.set(False)
    out = _markdown_to_docx_bytes(markdown, title)
    saved = _save_bytes(out_name, out)
    saved["content_type"] = DOCX_MIME
    saved["source_file"] = str(path.name)
    saved["conversion"] = "file_to_docx"
    saved["template_used"] = bool(_DOCX_TEMPLATE_USED.get())
    return saved


# -----------------------------
# PDF: remove_pages_save + merge_save (new)
# -----------------------------
def _infer_last_pdf_page(pdf_bytes: bytes) -> int:
    if PdfReader is None:
        raise HTTPException(status_code=500, detail="pdf_page_count_unavailable")
    try:
        reader = PdfReader(io.BytesIO(pdf_bytes))
        count = len(reader.pages)
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"pdf_page_count_failed: {exc}")
    if count < 1:
        raise HTTPException(status_code=400, detail="pdf_has_no_pages")
    return count


@app.post("/pdf/remove_pages_save", operation_id="pdf_remove_pages_save")
def pdf_remove_pages_save(payload: PdfRemovePagesSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)

    path = _resolve_path(payload.file_path, allowed_exts=(".pdf",), require_exact=True)
    data = _read_file(path)
    remove_pages = payload.remove_pages or payload.pages_to_remove or []
    if not remove_pages or payload.remove_last_page:
        remove_pages = [_infer_last_pdf_page(data)]

    files = {"file": (path.name, data, PDF_MIME)}
    form = {"remove_pages_json": json.dumps(remove_pages)}

    r = requests.post(
        f"{DOC_WORKER_URL}/pdf/remove_pages",
        headers=_worker_headers(),
        files=files,
        data=form,
        timeout=180,
    )
    _raise_for_worker_response(r)

    out_name = f"edited_{path.name}"
    saved = _save_bytes(out_name, r.content)
    saved["content_type"] = PDF_MIME
    saved["source_file"] = str(path.name)
    saved["removed_pages"] = remove_pages
    return saved


@app.post("/pdf/merge_save", operation_id="pdf_merge_save")
def pdf_merge_save(payload: PdfMergeSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)

    if not payload.file_paths or len(payload.file_paths) < 2:
        raise HTTPException(status_code=400, detail="provide_at_least_2_pdfs")

    paths: list[Path] = []
    for fp in payload.file_paths:
        _reject_wildcards(fp or "")
        p = _resolve_path(
            fp,
            allowed_exts=(".pdf",),
            require_exact=True,
            preferred_exact_paths=payload.attachment_exact_paths,
        )
        paths.append(p)

    mfiles = []
    for p in paths:
        b = _read_file(p)
        mfiles.append(("files", (p.name, b, PDF_MIME)))

    r = requests.post(
        f"{DOC_WORKER_URL}/pdf/merge",
        headers=_worker_headers(),
        files=mfiles,
        data={"debug_meta": "false"},
        timeout=300,
    )
    _raise_for_worker_response(r)

    out_name = payload.output_name or "merged.pdf"
    if not out_name.lower().endswith(".pdf"):
        out_name = f"{out_name}.pdf"

    saved = _save_bytes(out_name, r.content)
    saved["content_type"] = PDF_MIME
    saved["source_files"] = [p.name for p in paths]
    return saved


# -----------------------------
# Bundle to Markdown (existing + strict handling)
# -----------------------------
@app.post("/bundle/to_md", include_in_schema=False)
def bundle_to_md(payload: BundleToMdRequest, x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)

    if REQUIRE_EXACT_FILE_PATH and not payload.file_paths:
        raise HTTPException(status_code=400, detail="file_paths_required_in_strict_mode")

    title = (payload.title or "Masterkontext").strip() or "Masterkontext"

    paths: list[Path] = []
    if payload.file_paths:
        for fp in payload.file_paths:
            _reject_wildcards(fp or "")
            p = _resolve_path(fp, allowed_exts=(".docx", ".pdf", ".md", ".txt", ".xlsx", ".csv"))
            paths.append(p)
    else:
        # backwards-compatible best-effort
        for ext in ("docx", "pdf", "md", "txt", "xlsx", "csv"):
            try:
                paths.append(_pick_latest_file_by_exts((f".{ext}",)))
            except Exception:
                pass

    if not paths:
        raise HTTPException(status_code=400, detail="No input files found/provided.")

    mfiles = []
    for p in paths:
        b = _read_file(p)
        mfiles.append(("files", (p.name, b, "application/octet-stream")))

    r = requests.post(
        f"{DOC_WORKER_URL}/bundle/to_md",
        headers=_worker_headers(),
        files=mfiles,
        data={"title": title, "mode": "rag_mastercontext"},
        timeout=300,
    )
    _raise_for_worker_response(r)

    return Response(
        content=r.content,
        media_type="text/markdown; charset=utf-8",
        headers={"Content-Disposition": "attachment; filename=masterkontext.md"},
    )


@app.post("/bundle/to_md_save", operation_id="bundle_to_md_save")
def bundle_to_md_save(payload: BundleToMdSaveRequest, x_api_key: Optional[str] = Header(default=None)):
    _require_api_key(x_api_key)

    title = (payload.title or "Masterkontext").strip() or "Masterkontext"

    paths: list[Path] = []
    if payload.file_paths:
        for fp in payload.file_paths:
            _reject_wildcards(fp or "")
            p = _resolve_path(
                fp,
                allowed_exts=(".docx", ".pdf", ".md", ".txt", ".xlsx", ".csv"),
                require_exact=True,
            )
            paths.append(p)

    if not paths:
        raise HTTPException(status_code=400, detail="No input files found/provided.")

    mfiles = []
    for p in paths:
        b = _read_file(p)
        mfiles.append(("files", (p.name, b, "application/octet-stream")))

    r = requests.post(
        f"{DOC_WORKER_URL}/bundle/to_md",
        headers=_worker_headers(),
        files=mfiles,
        data={"title": title, "mode": "rag_mastercontext"},
        timeout=300,
    )
    _raise_for_worker_response(r)

    safe_title = _sanitize_filename(title)
    out_name = safe_title if safe_title.lower().endswith(".md") else f"{safe_title}.md"

    saved = _save_bytes(out_name, r.content)
    saved["content_type"] = "text/markdown; charset=utf-8"
    saved["source_files"] = [p.name for p in paths]
    return saved
